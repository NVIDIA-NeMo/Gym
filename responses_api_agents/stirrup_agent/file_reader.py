# Copyright (c) 2025, NVIDIA CORPORATION.  All rights reserved.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
"""Extract text content and convert deliverable files for reward computation.

Reads common document formats (.docx, .pdf, .xlsx, .pptx, .txt, etc.)
and returns their combined text so the LLM judge can score actual content
rather than just the agent's finish-tool summary.

Also provides PDF conversion via LibreOffice headless for visual judging
with multimodal models (e.g., Gemini 3 Pro).
"""

from __future__ import annotations

import base64
import math
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

# judge_panel owns the canonical audio/video extension sets and the GDPVal
# resources server routes tasks to judges from them. Importing rather than
# redeclaring keeps detection and emission in lockstep: a local copy that missed
# an extension would send a routed deliverable to a capable judge and then drop
# it here with no warning. judge_panel is stdlib-only, so this costs nothing.
from resources_servers.gdpval.judge_panel import AUDIO_EXTS, VIDEO_EXTS
from resources_servers.gdpval.preconvert import (
    AttachmentBudget,
    extract_xlsx_structured_text,
    resolve_pdf_provenance,
    roundtrip_ooxml_copy,
)


MAX_TOTAL_CHARS = 20_000


def _bounded_text(text: str, cap: int, marker: str = "\n[...truncated]") -> str:
    """Trim *text* including its marker so the returned string is at most *cap*."""

    if cap <= 0:
        return ""
    if len(text) <= cap:
        return text
    if cap <= len(marker):
        return marker[-cap:]
    return text[: cap - len(marker)] + marker


# Run state the agent writes into the SAME directory as its deliverables. These
# are not model output, and `.json` is in TEXT_EXTS, so without this the judge is
# shown the agent's own transcript and grades it as work product. Single source:
# every judging path imports it from here.
IGNORE_FILES = frozenset(
    {
        "finish_params.json",
        "history.json",
        "history.pkl",
        "inprogress_history.json",
        "metadata.json",
        "log.txt",
        "reference_files",
    }
)


def is_deliverable(path: Path) -> bool:
    """True if *path* is agent-produced output rather than run state."""
    return path.is_file() and path.name not in IGNORE_FILES


def read_deliverable_files(output_dir: str) -> str:
    """Read text from all deliverable files in *output_dir*.

    Returns a single string with sections like::

        === report.docx ===
        <extracted text>

        === data.xlsx ===
        <extracted text>

    Truncated to ~20k chars to stay within judge context limits.
    """
    output_path = Path(output_dir)
    if not output_path.is_dir():
        return ""

    files = sorted(output_path.iterdir())
    if not files:
        return ""

    parts: list[str] = []
    used = 0

    for fpath in files:
        if not is_deliverable(fpath):
            continue

        ext = fpath.suffix.lower()
        prefix = "\n\n" if parts else ""
        header = f"=== {fpath.name} ===\n"
        remaining = MAX_TOTAL_CHARS - used - len(prefix) - len(header)
        if remaining <= 0:
            if parts:
                parts[-1] = _bounded_text(parts[-1] + "x", len(parts[-1]))
            break
        try:
            text = _extract_text(fpath, ext, max_chars=remaining)
        except Exception as exc:
            text = _bounded_text(f"[Error reading {fpath.name}: {exc}]", remaining)

        if not text:
            continue

        section = header + _bounded_text(text, remaining)
        parts.append(section)
        used += len(prefix) + len(section)
        if used >= MAX_TOTAL_CHARS:
            break

    return _bounded_text("\n\n".join(parts), MAX_TOTAL_CHARS)


def _extract_text(fpath: Path, ext: str, max_chars: int = MAX_TOTAL_CHARS) -> str:
    """Dispatch to the right extractor based on file extension."""
    # TEXT_EXTS, not a second shorter list: a .ts/.py/.yaml deliverable was source
    # on the block path and "[Binary file: ...]" here.
    if ext in TEXT_EXTS:
        return _read_text(fpath, max_chars)
    elif ext == ".docx":
        return _read_docx(fpath, max_chars)
    elif ext == ".pdf":
        return _read_pdf(fpath, max_chars)
    elif ext == ".xlsx":
        return _read_xlsx(fpath, max_chars)
    elif ext == ".pptx":
        return _read_pptx(fpath, max_chars)
    else:
        size = os.path.getsize(fpath)
        return f"[Binary file: {fpath.name}, {size} bytes]"


def _read_text(fpath: Path, max_chars: int = MAX_TOTAL_CHARS) -> str:
    if max_chars <= 0:
        return ""
    with fpath.open("r", encoding="utf-8", errors="replace") as stream:
        text = stream.read(max_chars + 1)
    return _bounded_text(text, max_chars).strip()


def _read_docx(fpath: Path, max_chars: int = MAX_TOTAL_CHARS) -> str:
    from docx import Document

    doc = Document(str(fpath))
    parts: list[str] = []
    used = 0
    truncated = False
    for paragraph in doc.paragraphs:
        value = paragraph.text
        if not value.strip():
            continue
        separator = 1 if parts else 0
        remaining = max_chars - used - separator
        if remaining <= 0:
            truncated = True
            break
        parts.append(value[:remaining])
        used += separator + min(len(value), remaining)
        if len(value) > remaining:
            truncated = True
            break
    return _bounded_text("\n".join(parts) + ("x" if truncated else ""), max_chars)


def _read_pdf(fpath: Path, max_chars: int = MAX_TOTAL_CHARS) -> str:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer

    parts: list[str] = []
    used = 0
    truncated = False
    for page in extract_pages(str(fpath)):
        for element in page:
            if not isinstance(element, LTTextContainer):
                continue
            value = element.get_text()
            remaining = max_chars - used
            if remaining <= 0:
                truncated = True
                break
            parts.append(value[:remaining])
            used += min(len(value), remaining)
            if len(value) > remaining:
                truncated = True
                break
        if truncated:
            break
    text = "".join(parts).strip()
    return _bounded_text(text + ("x" if truncated else ""), max_chars)


def _read_xlsx(fpath: Path, max_chars: int = MAX_TOTAL_CHARS) -> str:
    return extract_xlsx_structured_text(fpath, max_chars=max_chars)


def _read_pptx(fpath: Path, max_chars: int = MAX_TOTAL_CHARS) -> str:
    from pptx import Presentation

    prs = Presentation(str(fpath))
    parts: list[str] = []
    used = 0
    truncated = False
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            value = f"Slide {i}:\n" + "\n".join(texts)
            separator = 2 if parts else 0
            remaining = max_chars - used - separator
            if remaining <= 0:
                truncated = True
                break
            parts.append(value[:remaining])
            used += separator + min(len(value), remaining)
            if len(value) > remaining:
                truncated = True
                break
    return _bounded_text("\n\n".join(parts) + ("x" if truncated else ""), max_chars)


# ---------------------------------------------------------------------------
# PDF conversion for visual judging (Gemini 3 Pro)
# ---------------------------------------------------------------------------

OOXML_EXTS = {".docx", ".pptx", ".xlsx"}
# LibreOffice renders these like OOXML; preconvert.py already converts them.
LEGACY_OFFICE_EXTS = {".doc", ".ppt", ".xls"}
OFFICE_EXTS = OOXML_EXTS | LEGACY_OFFICE_EXTS
TEXT_EXTS = (
    {".txt", ".md", ".rst", ".csv", ".tsv", ".json", ".jsonl", ".xml", ".html", ".htm", ".svg"}
    | {".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".properties", ".env"}
    | {".py", ".sh", ".bash", ".zsh", ".ps1", ".bat", ".sql", ".r", ".tex", ".ipynb"}
    | {".ts", ".tsx", ".js", ".jsx", ".mjs", ".cjs", ".vue", ".css", ".scss", ".less"}
    | {".java", ".go", ".rs", ".c", ".h", ".cpp", ".hpp", ".cc", ".cs", ".kt", ".swift"}
    | {".rb", ".php", ".pl", ".lua", ".scala", ".jl", ".m", ".log", ".diff", ".patch"}
)
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".heic", ".heif"}

# Derived, so it cannot drift from the branches below.
HANDLED_EXTS = TEXT_EXTS | OFFICE_EXTS | IMAGE_EXTS | AUDIO_EXTS | VIDEO_EXTS | {".pdf"}

TEXT_SNIFF_BYTES = 8192
# Per-file ceiling, and the aggregate across all text blocks in one request.
# Upstream has no cap at all: a big generated log produced a 720k-token request
# that the judge rejected outright, leaving the task silently unjudged.
MAX_TEXT_BLOCK_CHARS = 200_000
MAX_TOTAL_TEXT_BLOCK_CHARS = 400_000
# Binary payload ceilings for the rubric visual request. The encoded limit is
# the dominant wire-size term and leaves ample room below a 500 MB body ceiling
# for the bounded text above plus JSON framing.
MAX_FILE_ATTACHMENT_BYTES = 250 * 1024 * 1024
MAX_TOTAL_RAW_ATTACHMENT_BYTES = 300 * 1024 * 1024
MAX_TOTAL_ENCODED_ATTACHMENT_CHARS = 400 * 1024 * 1024
MAX_RASTER_PAGE_PIXELS = 40_000_000
# Header and truncation notice per block; reserved before allotting because
# it is tokens too.
TEXT_BLOCK_OVERHEAD_CHARS = 120
MAX_ARCHIVE_ENTRIES = 200
MAX_ARCHIVE_TEXT_CHARS = 120_000
MAX_ARCHIVE_MEMBER_BYTES = 2_000_000

# Document formats that happen to BE zip containers. Without this they open
# cleanly as archives and the judge is handed a listing of OOXML/ODF internals
# (`word/document.xml`, `[Content_Types].xml`) instead of the document.
ARCHIVE_DOC_EXTS = {
    ".xlsm",
    ".docm",
    ".pptm",
    ".xlsb",
    ".odt",
    ".ods",
    ".odp",
    ".odg",
    ".epub",
    ".jar",
    ".whl",
    ".apk",
    ".ipa",
}
# Audio/video (AUDIO_EXTS / VIDEO_EXTS, imported above) are only passed through
# when the judge is AV-capable (e.g. MiniMax M3); image-only judges can't decode
# them, so they are otherwise skipped. Every extension in those sets needs an
# entry here, or it routes to a capable judge and is then dropped silently.
MIME_TYPES = {
    ".pdf": "application/pdf",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".heic": "image/heic",
    ".heif": "image/heif",
    ".mp3": "audio/mpeg",
    ".wav": "audio/wav",
    ".m4a": "audio/mp4",
    ".flac": "audio/flac",
    ".ogg": "audio/ogg",
    ".oga": "audio/ogg",
    ".opus": "audio/opus",
    ".wma": "audio/x-ms-wma",
    ".aiff": "audio/aiff",
    ".aif": "audio/aiff",
    ".aac": "audio/aac",
    ".mp4": "video/mp4",
    ".m4v": "video/x-m4v",
    ".mov": "video/quicktime",
    ".webm": "video/webm",
    ".mkv": "video/x-matroska",
    ".avi": "video/x-msvideo",
    ".wmv": "video/x-ms-wmv",
    ".flv": "video/x-flv",
    ".mpeg": "video/mpeg",
    ".mpg": "video/mpeg",
    ".3gp": "video/3gpp",
}


def _human_size(n: float) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if n < 1024 or unit == "GB":
            return f"{n:.0f} B" if unit == "B" else f"{n:.1f} {unit}"
        n /= 1024.0
    return f"{n:.1f} GB"


def _attachment_omission(name: str, size: int, reason: str) -> dict[str, Any]:
    return {
        "type": "text",
        "text": f"\n{name}: [attachment omitted, {_human_size(size)} ({size:,} bytes): {reason}]",
    }


def _reserve_attachment(path: Path, budget: AttachmentBudget) -> tuple[int, dict[str, Any] | None]:
    """Reserve an attachment from metadata before reading or encoding it."""

    try:
        size = path.stat().st_size
    except OSError:
        return 0, {"type": "text", "text": f"\n{path.name}: [attachment unavailable]"}
    if size > MAX_FILE_ATTACHMENT_BYTES:
        return size, _attachment_omission(path.name, size, "per-file judge payload limit")
    if not budget.reserve(size):
        return size, _attachment_omission(path.name, size, "aggregate judge payload budget")
    return size, None


def _check_render_source(path: Path) -> tuple[int, dict[str, Any] | None]:
    """Reject a PDF source before loading it; rendered PNGs reserve separately."""

    try:
        size = path.stat().st_size
    except OSError:
        return 0, {"type": "text", "text": f"\n{path.name}: [attachment unavailable]"}
    if size > MAX_FILE_ATTACHMENT_BYTES:
        return size, _attachment_omission(path.name, size, "per-file judge payload limit")
    return size, None


def _sniffs_as_text(fpath: Path) -> bool:
    """True if an unknown-extension file looks like text (Makefile, Dockerfile, ...)."""
    try:
        with fpath.open("rb") as fh:
            chunk = fh.read(TEXT_SNIFF_BYTES)
    except OSError:
        return False
    if not chunk or b"\x00" in chunk:
        return False
    text = chunk.decode("utf-8", errors="replace")
    ctrl = sum(1 for ch in text if ord(ch) < 32 and ch not in "\t\n\r\f\v")
    if (text.count("�") + ctrl) / len(text) < 0.01:
        return True
    # cp1252/latin-1 is still text; only control characters make it binary.
    return ctrl / len(text) < 0.01


def _archive_manifest(fpath: Path) -> str | None:
    """Member listing plus text-member contents, or None if not a readable archive.

    Rubrics ask both whether a bundle contains X and whether what is inside is
    correct, so names alone cannot answer them.
    """
    import zipfile

    # Must come first: these open cleanly as zips and would otherwise be
    # summarised as their own XML internals.
    if fpath.suffix.lower() in ARCHIVE_DOC_EXTS:
        return None

    try:
        with zipfile.ZipFile(fpath) as zf:
            infos = zf.infolist()
            lines = [f"  {i.filename} ({i.file_size:,} bytes)" for i in infos[:MAX_ARCHIVE_ENTRIES]]
            if len(infos) > MAX_ARCHIVE_ENTRIES:
                lines.append(f"  ... and {len(infos) - MAX_ARCHIVE_ENTRIES:,} more entries")

            # Read IN MEMORY, never extracted: nothing is written to disk, so
            # there is no zip-slip to defend against. Bounded by per-member size,
            # aggregate characters, and the caller's own text budget.
            bodies: list[str] = []
            spent = 0
            for info in infos:
                if info.is_dir() or info.file_size > MAX_ARCHIVE_MEMBER_BYTES:
                    continue
                if Path(info.filename).suffix.lower() not in TEXT_EXTS:
                    continue
                if spent >= MAX_ARCHIVE_TEXT_CHARS:
                    bodies.append("[archive text budget exhausted; remaining members listed above but not shown]")
                    break
                try:
                    raw = zf.read(info)
                except Exception:
                    continue
                text = raw.decode("utf-8", errors="replace").strip()
                if not text:
                    continue
                take = text[: MAX_ARCHIVE_TEXT_CHARS - spent]
                spent += len(take)
                truncated = "\n[...member truncated]" if len(take) < len(text) else ""
                bodies.append(f"--- {info.filename} ---\n{take}{truncated}")
    except Exception:
        return None

    out = f"{len(infos)} member(s):\n" + "\n".join(lines)
    if bodies:
        out += "\n\n" + "\n\n".join(bodies)
    return out


def _fair_text_allowances(sizes: list[int], total_budget: int, per_file_cap: int) -> list[int]:
    """Max-min fair split of *total_budget* across text deliverables of *sizes*.

    Smallest first: each file may take an equal share of what is left divided by
    how many still need one. Files under their share take only what they need and
    hand back the surplus, so a small real deliverable is never starved by a large
    one that happens to sort earlier -- renaming a file must not change the score.
    Returned in the order *sizes* was given.
    """
    allowances = [0] * len(sizes)
    remaining = total_budget
    order = sorted(range(len(sizes)), key=lambda i: sizes[i])
    for position, idx in enumerate(order):
        still_to_serve = len(order) - position
        share = remaining // still_to_serve
        take = max(0, min(sizes[idx], per_file_cap, share))
        allowances[idx] = take
        remaining -= take
    return allowances


_libreoffice_unavailable_warned = False


def _warn_libreoffice_unavailable(exc: OSError) -> None:
    """Warn once per process that Office deliverables cannot be rendered.

    Per file this would be noise; never saying it at all is worse, because the
    silent consequence is that every Office deliverable is judged from extracted
    text instead of its rendered pages.
    """
    global _libreoffice_unavailable_warned
    if _libreoffice_unavailable_warned:
        return
    _libreoffice_unavailable_warned = True
    print(
        f"[file_reader] LibreOffice is unavailable ({exc}); Office deliverables will fall back to text "
        "extraction. Preconvert them to PDF for full-fidelity judging.",
        flush=True,
    )


def _convert_office_to_pdf(fpath: Path, out_dir: Path | None = None) -> Path | None:
    """Convert a .docx/.xlsx/.pptx file to PDF using LibreOffice headless.

    With *out_dir* the PDF is written there instead of beside *fpath*, so a
    judging pass never writes into the directory it is reading. The in-place
    default is kept because ``preconvert.py`` relies on it.

    Returns the path to the generated PDF, or None on failure.
    Uses a unique user profile to avoid lock conflicts in concurrent workers.

    Whitespace in the input filename makes LibreOffice's batch-convert mode
    silently drop the file (the URI it builds isn't percent-encoded), so we
    stage the input to a tempdir with a sanitized basename and move the PDF
    back to the original location.
    """
    dest_dir = out_dir if out_dir is not None else fpath.parent
    out_pdf = dest_dir / (fpath.stem + ".pdf")
    temp_dirs: list[Path] = []
    profile_dirs: list[Path] = []

    try:
        stage_dir: Path | None = None
        input_path = fpath
        has_whitespace = any(c.isspace() for c in fpath.name)
        if has_whitespace:
            stage_dir = Path(tempfile.mkdtemp(prefix="lo-stage-"))
            temp_dirs.append(stage_dir)
            safe_name = re.sub(r"\s+", "_", fpath.stem) + fpath.suffix
            input_path = stage_dir / safe_name
            shutil.copy2(fpath, input_path)
            first_outdir = stage_dir
        else:
            first_outdir = dest_dir

        def _invoke(source: Path, output: Path) -> tuple[subprocess.CompletedProcess[str], Path]:
            profile = Path(tempfile.mkdtemp(prefix="lo-profile-"))
            profile_dirs.append(profile)
            command = [
                "libreoffice",
                "--headless",
                "--nologo",
                "--nolockcheck",
                "--nodefault",
                "--norestore",
                f"-env:UserInstallation=file://{profile.as_posix()}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output),
                str(source),
            ]
            completed = subprocess.run(
                command,
                check=False,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=120,
            )
            return completed, output / (source.stem + ".pdf")

        first_result, first_pdf = _invoke(input_path, first_outdir)
        if first_pdf.exists():
            if first_pdf != out_pdf:
                shutil.move(str(first_pdf), str(out_pdf))
            return out_pdf

        retry_detail = ""
        if fpath.suffix.lower() in OOXML_EXTS:
            retry_dir = Path(tempfile.mkdtemp(prefix="gdpval-roundtrip-"))
            temp_dirs.append(retry_dir)
            retry_input = retry_dir / (re.sub(r"\s+", "_", fpath.stem) + fpath.suffix)
            try:
                roundtrip_ooxml_copy(fpath, retry_input)
                retry_result, retry_pdf = _invoke(retry_input, retry_dir)
                if retry_pdf.exists():
                    shutil.move(str(retry_pdf), str(out_pdf))
                    return out_pdf
                retry_detail = f"; round-trip retry rc={retry_result.returncode}: {retry_result.stderr[:200]}"
            except Exception as exc:
                retry_detail = f"; round-trip retry failed: {exc!r}"

        print(
            f"[file_reader] LibreOffice conversion failed for {fpath.name} "
            f"(rc={first_result.returncode}): {first_result.stderr[:200]}{retry_detail}",
            flush=True,
        )
        return None
    except subprocess.TimeoutExpired:
        print(f"[file_reader] LibreOffice conversion timed out for {fpath.name}", flush=True)
        return None
    except OSError as exc:
        # Usually LibreOffice simply isn't installed on this host. Returning None
        # degrades to the caller's text-extraction fallback; letting it propagate
        # would replace the entire deliverable with an "[Error: ...]" block, which a
        # judge reads as a missing artifact and scores near zero. That is not
        # hypothetical - it once manufactured a spurious "much worse" verdict.
        _warn_libreoffice_unavailable(exc)
        return None
    finally:
        for directory in profile_dirs + temp_dirs:
            shutil.rmtree(directory, ignore_errors=True)


def _pdf_bytes_to_image_text_blocks(
    pdf_bytes: bytes,
    *,
    render_dpi: int,
    max_pages: int,
    include_text: bool,
    attachment_budget: AttachmentBudget | None = None,
) -> list[dict[str, Any]] | None:
    """Rasterize one page at a time and reserve each PNG before base64."""

    try:
        import fitz
    except ImportError:
        return None

    budget = attachment_budget or AttachmentBudget(
        MAX_TOTAL_RAW_ATTACHMENT_BYTES,
        MAX_TOTAL_ENCODED_ATTACHMENT_CHARS,
    )
    try:
        document = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception:
        return None

    images: list[dict[str, Any]] = []
    text_parts: list[str] = []
    text_remaining = 20_000
    text_truncated = False
    try:
        total_pages = document.page_count
        page_limit = min(total_pages, max(0, max_pages))
        zoom = render_dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)
        for page_index in range(page_limit):
            try:
                page = document.load_page(page_index)
                if include_text and text_remaining > 0:
                    page_text = (page.get_text("text") or "").strip()
                    if page_text:
                        separator = 2 if text_parts else 0
                        available = max(0, text_remaining - separator)
                        take = min(len(page_text), available)
                        if take:
                            text_parts.append(page_text[:take])
                            text_remaining -= separator + take
                        if take < len(page_text):
                            text_truncated = True

                rect = page.rect
                width = math.ceil(rect.width * zoom)
                height = math.ceil(rect.height * zoom)
                if width * height > MAX_RASTER_PAGE_PIXELS:
                    images.append(
                        {
                            "type": "text",
                            "text": f"[page {page_index + 1} omitted: raster dimensions too large]",
                        }
                    )
                    continue
                pixmap = page.get_pixmap(matrix=matrix, alpha=False)
                png = pixmap.tobytes("png")
            except Exception:
                continue
            if not budget.reserve(len(png)):
                images.append(
                    {
                        "type": "text",
                        "text": "[remaining rendered pages omitted: aggregate judge payload budget]",
                    }
                )
                break
            b64 = base64.b64encode(png).decode("ascii")
            images.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}})

        if total_pages > page_limit:
            images.append({"type": "text", "text": f"[truncated: rendered {page_limit} of {total_pages} pages]"})
    finally:
        document.close()

    blocks: list[dict[str, Any]] = []
    if text_parts:
        text = "\n\n".join(text_parts)
        if text_truncated:
            text = _bounded_text(text + "x", 20_000, "\n[...text truncated]")
        blocks.append({"type": "text", "text": f"[extracted text]\n{text}"})
    blocks.extend(images)
    return blocks or None


def _av_block(mime: str, data: bytes, *, ext: str, file_type: str, openai_native: bool) -> dict[str, Any]:
    """Build an audio/video content block in the judge's dialect.

    Delegates to :mod:`resources_servers.gdpval.media_conversion` (imported
    lazily so this module stays importable where that package isn't on the
    path). Falls back to the legacy ``image_url`` data URL if the helper is
    unavailable — correct for frontier judges, and a benign degradation for
    self-hosted ones (which is only reachable when the gdpval package is
    present anyway).
    """
    try:
        from resources_servers.gdpval.media_conversion import audio_video_block

        return audio_video_block(mime, data, ext=ext, file_type=file_type, openai_native=openai_native)
    except ImportError:
        b64 = base64.b64encode(data).decode("ascii")
        return {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}}


def _enforce_text_block_budget(blocks: list[dict[str, Any]], cap: int) -> list[dict[str, Any]]:
    """Apply one hard character ceiling to every text producer in a request.

    The upstream builders include text that does not flow through ``_text_block``
    (PDF extraction, Office headers, AV stubs, archive manifests, and exception
    messages). This final pass is therefore the source of truth for the cap.
    Binary/media blocks are retained unchanged.
    """

    result: list[dict[str, Any]] = []
    remaining = max(0, cap)
    omitted = False
    last_text_index: int | None = None
    marker = "\n[...additional text blocks omitted: aggregate text budget exhausted]"

    for block in blocks:
        if block.get("type") != "text":
            result.append(block)
            continue
        text = str(block.get("text", ""))
        if omitted or remaining <= 0:
            omitted = omitted or bool(text)
            continue
        if len(text) > remaining:
            text = _bounded_text(text, remaining, marker)
            omitted = True
        copied = dict(block)
        copied["text"] = text
        result.append(copied)
        last_text_index = len(result) - 1
        remaining -= len(text)

    if omitted and last_text_index is not None and marker not in result[last_text_index]["text"]:
        previous = result[last_text_index]["text"]
        result[last_text_index]["text"] = _bounded_text(previous + marker, len(previous), marker)
    return result


def convert_deliverables_to_content_blocks(
    output_dir: str,
    *,
    media_mode: str = "native_pdf",
    render_dpi: int = 150,
    max_pages: int = 50,
    include_text: bool = True,
    audio_capable: bool = False,
    video_capable: bool = False,
) -> list[dict[str, Any]]:
    """Convert deliverable files to OpenAI-compatible content blocks for multimodal judging.

    Returns a list of content blocks suitable for the ``content`` field of an
    OpenAI chat message. Each file becomes either:

    - A text block (for .txt/.md/.csv etc.)
    - An image_url block with base64 data URL (for PDFs, images, converted Office docs)

    Office documents (.docx/.xlsx/.pptx) are first converted to PDF via
    LibreOffice headless so the judge sees rendered formatting/tables/charts.

    *media_mode* selects how PDFs/Office docs are presented: ``"native_pdf"``
    (default) sends them as ``application/pdf`` data URLs for frontier judges;
    ``"images_and_text"`` rasterizes each page to a PNG block and attaches the
    extracted text, for image-only local VLM judges (e.g. a gym-spawned Kimi
    K2.6). *render_dpi*, *max_pages*, and *include_text* tune that rendering.

    *audio_capable* / *video_capable* forward audio / video files (respectively)
    to judges that read that modality — tracked SEPARATELY because MiniMax-M3
    reads video but not audio (a video-only judge keeps video, skips audio). An
    unreadable modality's file is still announced by name and size rather than
    dropped. The AV block dialect follows
    *media_mode*: ``native_pdf`` (frontier judges, e.g. Gemini) uses an
    ``image_url`` data URL, while ``images_and_text`` (self-hosted vLLM judges)
    uses the standard ``video_url`` / ``input_audio`` content types vLLM routes to
    the media tower.
    """
    output_path = Path(output_dir)
    if not output_path.is_dir():
        return []

    images_and_text = media_mode == "images_and_text"

    blocks: list[dict[str, Any]] = []
    scratch_dirs: list[Path] = []  # tempdirs holding PDFs this pass converted
    attachment_budget = AttachmentBudget(
        MAX_TOTAL_RAW_ATTACHMENT_BYTES,
        MAX_TOTAL_ENCODED_ATTACHMENT_CHARS,
    )

    entries = sorted(output_path.iterdir())
    # Shared provenance rules keep the comparison and rubric paths in lockstep:
    # prefer source.ext.pdf, consume derived renders, and quarantine an old
    # source.pdf when multiple Office files share the stem.
    provenance = resolve_pdf_provenance(entry for entry in entries if is_deliverable(entry))
    consumed_pdfs = provenance.suppressed_pdfs

    # Which files the judge receives AS TEXT. Decided once and reused by both the
    # allotment and the dispatch loop: deciding it twice lets a file be emitted as
    # text while holding no allowance, i.e. truncated to nothing.
    texty = {
        p
        for p in entries
        if is_deliverable(p)
        and p not in consumed_pdfs
        and (p.suffix.lower() in TEXT_EXTS or (p.suffix.lower() not in HANDLED_EXTS and _sniffs_as_text(p)))
    }
    text_files = [p for p in entries if p in texty]
    body_budget = max(0, MAX_TOTAL_TEXT_BLOCK_CHARS - len(text_files) * TEXT_BLOCK_OVERHEAD_CHARS)
    allowance_by_name = dict(
        zip(
            (p.name for p in text_files),
            _fair_text_allowances([p.stat().st_size for p in text_files], body_budget, MAX_TEXT_BLOCK_CHARS),
        )
    )

    emitted_chars = 0
    omitted_names: list[str] = []

    def _text_block(header: str, body: str, allowance: int) -> dict[str, Any] | None:
        """A text block for *body*, trimmed to *allowance*.

        None once the aggregate budget is spent: with enough files even the headers
        alone would blow the request, so past that point they are summarised in one
        block rather than named individually.
        """
        nonlocal emitted_chars
        if emitted_chars >= MAX_TOTAL_TEXT_BLOCK_CHARS:
            return None
        cap = max(0, min(MAX_TEXT_BLOCK_CHARS, allowance))
        if len(body) > cap:
            shown = body[:cap] + f"\n[...truncated: {len(body) - cap:,} of {len(body):,} chars omitted]"
        else:
            shown = body
        block = {"type": "text", "text": f"\n{header}\n{shown}"}
        emitted_chars += len(block["text"])
        return block

    def _pdf_blocks(pdf_path: Path, header: str) -> list[dict[str, Any]]:
        """Emit a PDF natively or page-by-page under the shared binary budget."""

        size, source_omitted = _check_render_source(pdf_path)
        if source_omitted is not None:
            return [source_omitted]

        data: bytes | None = None
        if images_and_text:
            # The PDF is a bounded rasterization input, not request payload. Its
            # compressed source size does not predict whether a rendered PNG page
            # fits; the renderer reserves every actual page below.
            if not attachment_budget.can_fit(1):
                return [_attachment_omission(pdf_path.name, size, "aggregate judge payload budget")]
            data = pdf_path.read_bytes()
            rendered = _pdf_bytes_to_image_text_blocks(
                data,
                render_dpi=render_dpi,
                max_pages=max_pages,
                include_text=include_text,
                attachment_budget=attachment_budget,
            )
            if rendered is not None:
                return [{"type": "text", "text": header}, *rendered]

        if not attachment_budget.reserve(size):
            return [_attachment_omission(pdf_path.name, size, "aggregate judge payload budget")]
        if data is None:
            data = pdf_path.read_bytes()
        b64 = base64.b64encode(data).decode("ascii")
        return [
            {"type": "text", "text": header},
            {
                "type": "image_url",
                "image_url": {"url": f"data:application/pdf;base64,{b64}"},
            },
        ]

    for fpath in entries:
        if not is_deliverable(fpath) or fpath in consumed_pdfs:
            continue

        ext = fpath.suffix.lower()

        try:
            if fpath in texty:
                allowance = allowance_by_name.get(fpath.name, 0)
                if fpath.stat().st_size == 0:
                    block = _text_block(f"{fpath.name}:", "[present but EMPTY (0 bytes of content)]", 200)
                    blocks.append(block) if block else omitted_names.append(fpath.name)
                    continue
                if allowance <= 0:
                    omitted_names.append(fpath.name)
                    continue
                text = _read_text(fpath, allowance)
                if text:
                    block = _text_block(f"{fpath.name}:", text, allowance)
                else:
                    block = _text_block(f"{fpath.name}:", "[present but EMPTY (0 bytes of content)]", 200)
                blocks.append(block) if block else omitted_names.append(fpath.name)

            elif ext in OFFICE_EXTS:
                source_size = fpath.stat().st_size
                if source_size > MAX_FILE_ATTACHMENT_BYTES:
                    blocks.append(_attachment_omission(fpath.name, source_size, "per-file judge payload limit"))
                    continue
                # Prefer a PDF preconvert already rendered. Reconverting is
                # wasteful, needs LibreOffice on the judge host, and the cleanup
                # below would delete someone else's artifact.
                pdf_path = provenance.office_pdfs.get(fpath)
                if pdf_path is None:
                    scratch = Path(tempfile.mkdtemp(prefix="gdpval-render-"))
                    scratch_dirs.append(scratch)
                    pdf_path = _convert_office_to_pdf(fpath, out_dir=scratch)
                if pdf_path and pdf_path.exists():
                    header_kind = "rendered from PDF" if images_and_text else "converted to PDF"
                    blocks.extend(_pdf_blocks(pdf_path, f"\n{fpath.name} ({header_kind}):"))
                    if ext == ".xlsx":
                        sheet_text = extract_xlsx_structured_text(fpath, max_chars=MAX_TEXT_BLOCK_CHARS)
                        block = _text_block(
                            f"{fpath.name} (structured spreadsheet cells):",
                            sheet_text,
                            MAX_TEXT_BLOCK_CHARS,
                        )
                        blocks.append(block) if block else omitted_names.append(fpath.name)
                else:
                    # Fallback to text extraction
                    text = _extract_text(fpath, ext, max_chars=MAX_TEXT_BLOCK_CHARS)
                    if text:
                        block = _text_block(f"{fpath.name} (text fallback):", text, MAX_TEXT_BLOCK_CHARS)
                        blocks.append(block) if block else omitted_names.append(fpath.name)

            elif ext == ".pdf":
                blocks.extend(_pdf_blocks(fpath, f"\n{fpath.name}:"))

            elif ext in IMAGE_EXTS:
                mime = MIME_TYPES.get(ext, "image/png")
                _size, omitted = _reserve_attachment(fpath, attachment_budget)
                if omitted is not None:
                    blocks.append(omitted)
                    continue
                data = fpath.read_bytes()
                b64 = base64.b64encode(data).decode("ascii")
                blocks.append({"type": "text", "text": f"\n{fpath.name}:"})
                blocks.append(
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime};base64,{b64}"},
                    }
                )

            elif ext in AUDIO_EXTS or ext in VIDEO_EXTS:
                # Forward AV only to judges that read that SPECIFIC modality, gated
                # independently: video needs *video_capable*, audio *audio_capable*
                # — so MiniMax-M3 (video yes, audio no) keeps video and skips audio.
                # Among frontier judges only Gemini reads AV. Gemini (native_pdf)
                # accepts AV as an image_url data URL; self-hosted vLLM judges
                # (images_and_text) need the standard video_url / input_audio types,
                # which vLLM routes to the media tower.
                is_video = ext in VIDEO_EXTS
                file_type = "VIDEO" if is_video else "AUDIO"
                if video_capable if is_video else audio_capable:
                    mime = MIME_TYPES.get(ext, "application/octet-stream")
                    _size, omitted = _reserve_attachment(fpath, attachment_budget)
                    if omitted is not None:
                        blocks.append(omitted)
                        continue
                    data = fpath.read_bytes()
                    blocks.append({"type": "text", "text": f"\n{fpath.name}:"})
                    blocks.append(_av_block(mime, data, ext=ext, file_type=file_type, openai_native=images_and_text))
                else:
                    # Gating is correct; emitting nothing is not -- the judge
                    # then grades against a file it believes was never produced.
                    size = fpath.stat().st_size
                    blocks.append(
                        {
                            "type": "text",
                            "text": (
                                f"\n{fpath.name}: [{file_type} deliverable, {_human_size(size)} "
                                f"({size:,} bytes) — present on disk, but this judge cannot decode "
                                f"{file_type.lower()}. Do NOT treat it as missing or unproduced. Grade the "
                                f"criteria that do not require {'viewing' if is_video else 'listening'}; "
                                f"mark the rest unverifiable rather than unmet. "
                                f"Statements in other files about this file's contents are unverified "
                                f"assertions, not evidence: do not award credit for a property you "
                                f"cannot observe directly.]"
                            ),
                        }
                    )

            else:
                # Any allowlist lags reality; an unlisted extension used to emit
                # nothing, so the judge scored it as never produced.
                size = fpath.stat().st_size
                if size == 0:
                    body = "[present but EMPTY (0 bytes)]"
                elif _sniffs_as_text(fpath):
                    text = _read_text(fpath, MAX_TEXT_BLOCK_CHARS)
                    body = text if text else "[present but EMPTY (0 bytes of content)]"
                else:
                    manifest = _archive_manifest(fpath)
                    if manifest is not None:
                        body = f"[archive, {_human_size(size)} — NOT missing] {manifest}"
                    else:
                        body = (
                            f"[deliverable present, {_human_size(size)} ({size:,} bytes) — NOT missing; "
                            f"content not extractable in this format. Statements in other files about "
                            f"its contents are unverified assertions, not evidence.]"
                        )
                block = _text_block(f"{fpath.name}:", body, MAX_TEXT_BLOCK_CHARS)
                blocks.append(block) if block else omitted_names.append(fpath.name)
        except Exception as exc:
            blocks.append({"type": "text", "text": f"\n{fpath.name}: [Error: {exc}]"})

    if omitted_names:
        shown = ", ".join(omitted_names[:20])
        if len(omitted_names) > 20:
            shown += f", and {len(omitted_names) - 20} more"
        blocks.append(
            {
                "type": "text",
                "text": f"\n[{len(omitted_names)} text deliverable(s) omitted, budget exhausted: {shown}]",
            }
        )

    # Clean up only what this pass created; it all lives in tempdirs.
    for scratch in scratch_dirs:
        shutil.rmtree(scratch, ignore_errors=True)

    return _enforce_text_block_budget(blocks, MAX_TOTAL_TEXT_BLOCK_CHARS)
