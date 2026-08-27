# SPDX-FileCopyrightText: Copyright (c) 2026 NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0

"""Local extraction and rendering for APEX grading artifacts.

The document extraction behavior in this module is adapted from the Apache-2.0
Mercor Archipelago grader at commit 0cb5c476c219a9df637e0bd37fb86b2361f4ab89.
Only its offline artifact-processing layer is carried here; model invocation,
artifact-selection LLMs, cloud extractors, and scoring remain owned by NeMo Gym.
"""

from __future__ import annotations

import base64
import csv
import io
import json
import math
import re
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET


_TEXT_EXTENSIONS = {
    ".html",
    ".json",
    ".log",
    ".md",
    ".py",
    ".sh",
    ".sql",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}
_OFFICE_EXTENSIONS = {".docx", ".pptx", ".xls", ".xlsm", ".xlsx"}
_IMAGE_MIME_TYPES = {
    ".heic": "image/heic",
    ".heif": "image/heif",
    ".jpeg": "image/jpeg",
    ".jpg": "image/jpeg",
    ".png": "image/png",
    ".webp": "image/webp",
}
_WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_DOCUMENT_XML_PATH = "word/document.xml"
_COMMENTS_XML_PATH = "word/comments.xml"
_SECTION_HEADING_RE = re.compile(r"^(SECTION|ARTICLE)\s+[0-9IVXLC]+\b|^[0-9]+\.[0-9]+")
_MAX_VISUAL_PAGES = 10


@dataclass(frozen=True)
class SubArtifact:
    """One Archipelago-style slide, sheet, or page within an artifact."""

    index: int
    type: str
    title: str | None
    content: str
    images: list[dict[str, Any]] = field(default_factory=list)


@dataclass(frozen=True)
class ExtractedContent:
    """Structured extraction result with concatenated text for compatibility."""

    text: str
    images: list[dict[str, Any]] = field(default_factory=list)
    sub_artifacts: list[SubArtifact] = field(default_factory=list)


def extract_file_text(path: Path, *, document_converter_image: str | None = None) -> str:
    """Extract judge-readable text using Archipelago's offline format semantics."""
    return extract_file_content(path, document_converter_image=document_converter_image).text


def extract_file_content(path: Path, *, document_converter_image: str | None = None) -> ExtractedContent:
    """Extract text and Archipelago-compatible sub-artifact records."""
    extension = path.suffix.lower()
    if extension in _TEXT_EXTENSIONS:
        text = path.read_text(encoding="utf-8", errors="replace").strip()
        return ExtractedContent(text=text)
    if extension == ".csv":
        return ExtractedContent(text=_extract_csv(path))
    if extension == ".ipynb":
        return ExtractedContent(text=_extract_ipynb(path))
    if extension == ".docx":
        return ExtractedContent(text=_extract_docx(path))
    if extension == ".pdf":
        return _extract_pdf(path)
    if extension in {".xlsx", ".xlsm"}:
        return _extract_xlsx(path, document_converter_image=document_converter_image)
    if extension == ".xls":
        return _extract_xls(path)
    if extension == ".pptx":
        return _extract_pptx(path)
    return ExtractedContent(text=f"[Binary file: {path.name}, {path.stat().st_size} bytes]")


def _extract_csv(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            with path.open(encoding=encoding, newline="") as stream:
                return "\n".join("\t".join(row) for row in csv.reader(stream))
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _join_ipynb_source(source: Any) -> str:
    return "".join(str(line) for line in source) if isinstance(source, list) else str(source or "")


def _ipynb_text_outputs(outputs: Any) -> str:
    if not isinstance(outputs, list):
        return ""
    lines: list[str] = []
    for output in outputs:
        if not isinstance(output, dict):
            continue
        output_type = output.get("output_type")
        if output_type == "stream":
            lines.append(_join_ipynb_source(output.get("text")))
        elif output_type in {"execute_result", "display_data"}:
            data = output.get("data")
            if isinstance(data, dict) and "text/plain" in data:
                lines.append(_join_ipynb_source(data["text/plain"]))
        elif output_type == "error":
            traceback = output.get("traceback")
            if isinstance(traceback, list):
                lines.append("\n".join(str(line) for line in traceback))
    return "\n".join(line.strip() for line in lines if line.strip())


def _extract_ipynb(path: Path) -> str:
    data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    parts: list[str] = []
    for cell in data.get("cells", []) if isinstance(data, dict) else []:
        if not isinstance(cell, dict):
            continue
        source = _join_ipynb_source(cell.get("source")).rstrip()
        if cell.get("cell_type") == "markdown" and source:
            parts.append(source)
        elif cell.get("cell_type") == "code":
            if source:
                parts.append(f"```python\n{source}\n```")
            output = _ipynb_text_outputs(cell.get("outputs"))
            if output:
                parts.append(f"Output:\n{output}")
    return "\n\n".join(parts)


def _extract_text_from_pptx_shape(shape: Any) -> list[str]:
    """Recursively extract PowerPoint text, including grouped shapes and tables."""
    text_parts: list[str] = []
    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            text_parts.extend(_extract_text_from_pptx_shape(child_shape))
        return text_parts

    try:
        table = shape.table
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if row_texts:
                text_parts.append("\t".join(row_texts))
        return text_parts
    except (ValueError, AttributeError):
        pass

    text_frame_succeeded = False
    if hasattr(shape, "text_frame"):
        try:
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = "".join(run.text for run in paragraph.runs if run.text)
                if paragraph_text.strip():
                    text_parts.append(paragraph_text.strip())
                    text_frame_succeeded = True
        except Exception:
            pass
    if not text_frame_succeeded and hasattr(shape, "text") and shape.text:
        text = shape.text.strip()
        if text:
            text_parts.append(text)
    return text_parts


def _extract_pptx(path: Path) -> ExtractedContent:
    from pptx import Presentation

    presentation = Presentation(str(path))
    full_text_parts: list[str] = []
    sub_artifacts: list[SubArtifact] = []
    for slide_index, slide in enumerate(presentation.slides):
        slide_text_parts: list[str] = []
        slide_title: str | None = None
        for shape in slide.shapes:
            if slide_title is None:
                try:
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1 and shape.text:
                        slide_title = shape.text.strip()
                except Exception:
                    pass
            slide_text_parts.extend(_extract_text_from_pptx_shape(shape))
        if slide_title is None and slide_text_parts:
            slide_title = slide_text_parts[0][:100]
        slide_text = "\n".join(slide_text_parts)
        sub_artifacts.append(
            SubArtifact(
                index=slide_index,
                type="slide",
                title=slide_title or f"Slide {slide_index + 1}",
                content=slide_text,
            )
        )
        full_text_parts.append(f"=== Slide {slide_index + 1}: {slide_title or 'Untitled'} ===\n{slide_text}")
    return ExtractedContent(text="\n\n".join(full_text_parts), sub_artifacts=sub_artifacts)


def _extract_pdf(path: Path) -> ExtractedContent:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        reader.decrypt("")
    pages: list[str] = []
    sub_artifacts: list[SubArtifact] = []
    for page_index, page in enumerate(reader.pages):
        try:
            page_text = (page.extract_text() or "").strip()
        except Exception as exc:
            page_text = f"[page {page_index + 1} text extraction failed: {type(exc).__name__}]"
        sub_artifacts.append(
            SubArtifact(
                index=page_index,
                type="page",
                title=f"Page {page_index + 1}",
                content=page_text,
            )
        )
        pages.append(f"=== Page {page_index + 1} ===\n{page_text}")
    scanned_banner = ""
    if sub_artifacts and not any(page.content.strip() for page in sub_artifacts):
        scanned_banner = (
            "=== Extraction Note ===\n"
            "[no extractable text — the PDF is likely scanned/image-only and no OCR is available offline]\n\n"
        )
        sub_artifacts = [
            SubArtifact(
                index=page.index,
                type=page.type,
                title=page.title,
                content=scanned_banner + page.content,
                images=page.images,
            )
            for page in sub_artifacts
        ]
    return ExtractedContent(
        text=scanned_banner + "\n\n".join(pages),
        sub_artifacts=sub_artifacts,
    )


def _accepted_docx_bytes(raw_bytes: bytes) -> bytes:
    """Return DOCX bytes with revisions accepted, matching Archipelago text transformation."""
    try:
        with zipfile.ZipFile(io.BytesIO(raw_bytes), "r") as source:
            if _DOCUMENT_XML_PATH not in source.namelist():
                return raw_bytes
            from lxml import etree

            parser = etree.XMLParser(resolve_entities=False, no_network=True)
            tree = etree.fromstring(source.read(_DOCUMENT_XML_PATH), parser=parser)
            for local_name in ("del", "moveFrom"):
                for element in list(tree.iter(f"{{{_WORD_NS}}}{local_name}")):
                    if (parent := element.getparent()) is not None:
                        parent.remove(element)
            for local_name in ("ins", "moveTo"):
                for element in list(tree.iter(f"{{{_WORD_NS}}}{local_name}")):
                    parent = element.getparent()
                    if parent is None:
                        continue
                    index = list(parent).index(element)
                    for child in reversed(list(element)):
                        parent.insert(index, child)
                    parent.remove(element)
            modified_xml = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)
            destination = io.BytesIO()
            with zipfile.ZipFile(destination, "w", zipfile.ZIP_DEFLATED) as output:
                for item in source.infolist():
                    output.writestr(
                        item, modified_xml if item.filename == _DOCUMENT_XML_PATH else source.read(item.filename)
                    )
            return destination.getvalue()
    except Exception:
        return raw_bytes


def _docx_review_text(raw_bytes: bytes) -> str:
    """Expose tracked changes and anchored comments that python-docx omits."""
    try:
        from lxml import etree

        parser = etree.XMLParser(resolve_entities=False, no_network=True)
        with zipfile.ZipFile(io.BytesIO(raw_bytes), "r") as archive:
            if _DOCUMENT_XML_PATH not in archive.namelist():
                return ""
            document = etree.fromstring(archive.read(_DOCUMENT_XML_PATH), parser=parser)
            comments: dict[str, tuple[str, str]] = {}
            if _COMMENTS_XML_PATH in archive.namelist():
                comments_root = etree.fromstring(archive.read(_COMMENTS_XML_PATH), parser=parser)
                for comment in comments_root.iter(f"{{{_WORD_NS}}}comment"):
                    comment_id = comment.get(f"{{{_WORD_NS}}}id")
                    if comment_id is not None:
                        author = comment.get(f"{{{_WORD_NS}}}author") or "Unknown author"
                        text = "".join(node.text or "" for node in comment.iter(f"{{{_WORD_NS}}}t")).strip()
                        comments[comment_id] = (author, text)

        redlines_by_author: dict[str, list[str]] = {}
        anchors: dict[str, dict[str, Any]] = {}
        open_comment_ids: set[str] = set()
        current_section = ""
        paragraph_index = 0
        revision_tags = {"ins": "INSERTED", "del": "DELETED", "moveFrom": "MOVED FROM", "moveTo": "MOVED TO"}
        for paragraph in document.iter(f"{{{_WORD_NS}}}p"):
            paragraph_index += 1
            paragraph_text = "".join(node.text or "" for node in paragraph.iter(f"{{{_WORD_NS}}}t")).strip()
            if paragraph_text and _SECTION_HEADING_RE.match(paragraph_text):
                current_section = paragraph_text[:60].strip()
            location = current_section or f"Para {paragraph_index}"
            for element in paragraph.iter():
                local_name = element.tag.rsplit("}", 1)[-1]
                if local_name == "commentRangeStart":
                    comment_id = element.get(f"{{{_WORD_NS}}}id")
                    if comment_id is not None:
                        anchors[comment_id] = {"location": location, "parts": []}
                        open_comment_ids.add(comment_id)
                elif local_name == "commentRangeEnd":
                    comment_id = element.get(f"{{{_WORD_NS}}}id")
                    if comment_id is not None:
                        open_comment_ids.discard(comment_id)
                elif local_name == "t" and element.text:
                    for comment_id in open_comment_ids:
                        anchors[comment_id]["parts"].append(element.text)
            for revision_name, label in revision_tags.items():
                text_tag = "delText" if revision_name == "del" else "t"
                for revision in paragraph.iter(f"{{{_WORD_NS}}}{revision_name}"):
                    text = "".join(node.text or "" for node in revision.iter(f"{{{_WORD_NS}}}{text_tag}"))
                    if text:
                        author = revision.get(f"{{{_WORD_NS}}}author") or "Unknown author"
                        redlines_by_author.setdefault(author, []).append(f'  [{location}] {label}: "{text}"')

        parts: list[str] = []
        if redlines_by_author:
            parts.append("=== DOCUMENT REDLINES ===")
            for author, lines in redlines_by_author.items():
                parts.extend((f"\nAuthor: {author}", *lines))
        if comments:
            comment_lines = ["=== DOCUMENT COMMENTS ==="]
            for comment_id, (author, text) in comments.items():
                anchor = anchors.get(comment_id)
                if anchor:
                    anchored_text = "".join(anchor["parts"]).strip()
                    on = f' on "{anchored_text}"' if anchored_text else ""
                    comment_lines.append(f"\nAuthor: {author}\n  [{anchor['location']}]{on}: {text}")
                else:
                    comment_lines.append(f"\nAuthor: {author}\n  [no anchor]: {text}")
            parts.append("\n".join(comment_lines))
        return "\n\n".join(parts)
    except Exception:
        return ""


def _extract_docx(path: Path) -> str:
    from docx import Document

    raw_bytes = path.read_bytes()
    document = Document(io.BytesIO(_accepted_docx_bytes(raw_bytes)))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in getattr(document, "tables", []):
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    review_text = _docx_review_text(raw_bytes)
    if review_text:
        parts.append(review_text)
    return "\n".join(parts)


def _hidden_xlsx_sheets(path: Path) -> set[str]:
    hidden: set[str] = set()
    try:
        with zipfile.ZipFile(path, "r") as archive, archive.open("xl/workbook.xml") as stream:
            root = ET.parse(stream).getroot()
        namespace = {"main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        for sheet in root.findall(".//main:sheet", namespace):
            if sheet.get("state", "visible") in {"hidden", "veryHidden"} and sheet.get("name"):
                hidden.add(str(sheet.get("name")))
    except Exception:
        pass
    return hidden


def _resolve_document_command(executable: str, document_converter_image: str | None) -> tuple[str, ...] | None:
    if local_executable := shutil.which(executable):
        return (local_executable,)
    apptainer = shutil.which("apptainer")
    if apptainer and document_converter_image and Path(document_converter_image).is_file():
        return (apptainer, "exec", "--cleanenv", document_converter_image, executable)
    return None


def _xlsx_has_uncached_formulas(path: Path) -> bool:
    """Port Archipelago's formula/cache coordinate comparison."""
    from openpyxl import load_workbook

    try:
        formulas = load_workbook(path, data_only=False, read_only=True)
        try:
            formula_coordinates = {
                (sheet_name, cell.coordinate)
                for sheet_name in formulas.sheetnames
                for row in formulas[sheet_name].iter_rows()
                for cell in row
                if cell.data_type == "f" or (cell.value is not None and str(cell.value).startswith("="))
            }
        finally:
            formulas.close()
        if not formula_coordinates:
            return False

        values = load_workbook(path, data_only=True, read_only=True)
        try:
            return any(
                cell.value is None and (sheet_name, getattr(cell, "coordinate", None)) in formula_coordinates
                for sheet_name in values.sheetnames
                for row in values[sheet_name].iter_rows()
                for cell in row
            )
        finally:
            values.close()
    except Exception:
        return True


def _evaluate_excel_formulas_with_libreoffice(
    path: Path,
    *,
    document_converter_image: str | None,
) -> bytes | None:
    """Re-save a workbook through LibreOffice so formula cells gain cached values."""
    command = _resolve_document_command("libreoffice", document_converter_image) or _resolve_document_command(
        "soffice", document_converter_image
    )
    if command is None:
        return None

    input_path: Path | None = None
    output_dir: Path | None = None
    profile_dir: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=path.suffix, prefix="formula_eval_") as stream:
            stream.write(path.read_bytes())
            input_path = Path(stream.name)
        output_dir = Path(tempfile.mkdtemp(prefix="formula_eval_out_"))
        profile_dir = Path(tempfile.mkdtemp(prefix="libreoffice_profile_"))
        process = subprocess.run(
            [
                *command,
                "--headless",
                "--calc",
                f"-env:UserInstallation=file://{profile_dir}",
                "--convert-to",
                "xlsx",
                "--outdir",
                str(output_dir),
                str(input_path),
            ],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,
        )
        output_path = output_dir / f"{input_path.stem}.xlsx"
        return output_path.read_bytes() if process.returncode == 0 and output_path.is_file() else None
    except (OSError, subprocess.TimeoutExpired):
        return None
    finally:
        if input_path is not None:
            input_path.unlink(missing_ok=True)
        if output_dir is not None:
            shutil.rmtree(output_dir, ignore_errors=True)
        if profile_dir is not None:
            shutil.rmtree(profile_dir, ignore_errors=True)


def _has_charts_in_xlsx(path: Path) -> bool:
    from openpyxl import load_workbook

    workbook = None
    try:
        workbook = load_workbook(path, data_only=True)
        return any(bool(getattr(workbook[sheet_name], "_charts", None)) for sheet_name in workbook.sheetnames)
    except Exception:
        return False
    finally:
        if workbook is not None:
            workbook.close()


def _pdf_page_data_urls(
    pdf_path: Path,
    *,
    document_converter_image: str | None,
    page_indices: set[int] | None = None,
) -> list[tuple[int, str]]:
    page_numbers = sorted(index + 1 for index in page_indices)[:_MAX_VISUAL_PAGES] if page_indices else None
    command = _resolve_document_command("pdftoppm", document_converter_image)
    if command is None:
        from pdf2image import convert_from_path

        ranges = (
            [(page_number, page_number) for page_number in page_numbers] if page_numbers else [(1, _MAX_VISUAL_PAGES)]
        )
        results: list[tuple[int, str]] = []
        for first_page, last_page in ranges:
            pages = convert_from_path(str(pdf_path), dpi=150, first_page=first_page, last_page=last_page)
            for offset, page in enumerate(pages):
                buffer = io.BytesIO()
                page.save(buffer, format="PNG")
                encoded = base64.b64encode(buffer.getvalue()).decode("ascii")
                results.append((first_page + offset, f"data:image/png;base64,{encoded}"))
        return results

    output_dir = Path(tempfile.mkdtemp(prefix="apex-chart-pages-"))
    try:
        if page_numbers:
            selected: list[tuple[int, str]] = []
            for page_number in page_numbers:
                prefix = output_dir / f"selected-{page_number}"
                process = subprocess.run(
                    [
                        *command,
                        "-png",
                        "-r",
                        "150",
                        "-f",
                        str(page_number),
                        "-l",
                        str(page_number),
                        "-singlefile",
                        str(pdf_path),
                        str(prefix),
                    ],
                    check=False,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=120,
                )
                image_path = prefix.with_suffix(".png")
                if process.returncode == 0 and image_path.is_file():
                    encoded = base64.b64encode(image_path.read_bytes()).decode("ascii")
                    selected.append((page_number, f"data:image/png;base64,{encoded}"))
            return selected

        prefix = output_dir / "page"
        process = subprocess.run(
            [*command, "-png", "-r", "150", "-f", "1", "-l", str(_MAX_VISUAL_PAGES), str(pdf_path), str(prefix)],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,
        )
        if process.returncode != 0:
            return []
        results: list[tuple[int, str]] = []
        for page_number, image_path in enumerate(sorted(output_dir.glob("page-*.png")), 1):
            encoded = base64.b64encode(image_path.read_bytes()).decode("ascii")
            results.append((page_number, f"data:image/png;base64,{encoded}"))
        return results
    except (OSError, subprocess.TimeoutExpired):
        return []
    finally:
        shutil.rmtree(output_dir, ignore_errors=True)


def _pdf_to_chart_images(pdf_path: Path, *, document_converter_image: str | None) -> list[dict[str, Any]]:
    try:
        return [
            {
                "url": url,
                "placeholder": f"[CHART_{index}]",
                "type": "Chart",
                "caption": f"Chart from Excel (Page {page_number})",
            }
            for index, (page_number, url) in enumerate(
                _pdf_page_data_urls(pdf_path, document_converter_image=document_converter_image),
                1,
            )
        ]
    except Exception:
        return []


def _extract_xlsx(path: Path, *, document_converter_image: str | None = None) -> ExtractedContent:
    from openpyxl import load_workbook

    hidden_sheets = _hidden_xlsx_sheets(path)
    formulas: dict[tuple[str, str], str] = {}
    formula_workbook = load_workbook(path, data_only=False, read_only=True)
    try:
        for sheet_name in formula_workbook.sheetnames:
            if sheet_name in hidden_sheets:
                continue
            for row in formula_workbook[sheet_name].iter_rows():
                for cell in row:
                    if cell.data_type == "f" and isinstance(cell.value, str) and cell.value.startswith("="):
                        formulas[(sheet_name, cell.coordinate)] = cell.value
    finally:
        formula_workbook.close()

    unresolved = _xlsx_has_uncached_formulas(path)
    recalculated = (
        _evaluate_excel_formulas_with_libreoffice(path, document_converter_image=document_converter_image)
        if unresolved
        else None
    )
    if recalculated is not None:
        unresolved = False
    workbook_source: Path | io.BytesIO = io.BytesIO(recalculated) if recalculated is not None else path
    workbook = load_workbook(workbook_source, data_only=True, read_only=True)
    parts: list[str] = []
    sheet_values: list[tuple[int, str, str]] = []
    try:
        for sheet_index, sheet_name in enumerate(workbook.sheetnames):
            if sheet_name in hidden_sheets:
                continue
            rows: list[str] = []
            formula_lines: list[str] = []
            format_lines: list[str] = []
            for row in workbook[sheet_name].iter_rows():
                values: list[str] = []
                for cell in row:
                    coordinate = getattr(cell, "coordinate", "")
                    formula = formulas.get((sheet_name, coordinate))
                    if formula:
                        formula_lines.append(f"{coordinate}: {formula}")
                    if cell.value is None:
                        continue
                    values.append(str(cell.value))
                    if isinstance(cell.value, (int, float)) and not isinstance(cell.value, bool):
                        number_format = cell.number_format
                        if number_format and number_format != "General" and re.search(r"[$€£¥%]|#,##0", number_format):
                            format_lines.append(f"{coordinate}: [fmt: {number_format}]")
                if values:
                    rows.append("\t".join(values))
            sheet_text = f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows)
            if format_lines:
                sheet_text += "\n\n=== Number Formats ===\n" + "\n".join(format_lines)
            if formula_lines:
                sheet_text += "\n\n=== Formulas ===\n" + "\n".join(formula_lines)
            parts.append(sheet_text)
            sheet_values.append((sheet_index, sheet_name, sheet_text))
    finally:
        workbook.close()
    banner = ""
    if unresolved:
        banner = (
            "=== Extraction Note ===\n[formula cells could not be recalculated — LibreOffice unavailable; "
            "blanks may be uncomputed formulas, not missing values]\n\n"
        )
    sub_artifacts = [
        SubArtifact(index=index, type="sheet", title=title, content=banner + content)
        for index, title, content in sheet_values
    ]
    chart_images: list[dict[str, Any]] = []
    if _has_charts_in_xlsx(path):
        pdf_path = _convert_office_to_pdf(path, document_converter_image=document_converter_image)
        try:
            if pdf_path is not None:
                chart_images = _pdf_to_chart_images(
                    pdf_path,
                    document_converter_image=document_converter_image,
                )
        finally:
            if pdf_path is not None:
                shutil.rmtree(pdf_path.parent, ignore_errors=True)
    chart_text = ""
    if chart_images:
        chart_text = "\n\n=== Charts ===\n" + "\n".join(
            f"{image['placeholder']} - {image['caption']}" for image in chart_images
        )
    return ExtractedContent(
        text=banner + "\n\n".join(parts) + chart_text,
        images=chart_images,
        sub_artifacts=sub_artifacts,
    )


def _extract_xls(path: Path) -> ExtractedContent:
    import xlrd

    workbook = xlrd.open_workbook(str(path))
    parts: list[str] = []
    sub_artifacts: list[SubArtifact] = []
    for sheet_index in range(workbook.nsheets):
        sheet = workbook.sheet_by_index(sheet_index)
        if sheet.visibility != 0:
            continue
        rows: list[str] = []
        for row_index in range(sheet.nrows):
            values: list[str] = []
            for column_index in range(sheet.ncols):
                cell = sheet.cell(row_index, column_index)
                if cell.ctype == xlrd.XL_CELL_EMPTY:
                    continue
                if cell.ctype == xlrd.XL_CELL_NUMBER:
                    value = cell.value
                    if (
                        isinstance(value, float)
                        and not (math.isnan(value) or math.isinf(value))
                        and value == int(value)
                    ):
                        values.append(str(int(value)))
                    else:
                        values.append(str(value))
                elif cell.ctype == xlrd.XL_CELL_DATE:
                    date = xlrd.xldate_as_tuple(float(cell.value), workbook.datemode)
                    values.append(f"{date[0]}-{date[1]:02d}-{date[2]:02d}")
                elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                    values.append("TRUE" if cell.value else "FALSE")
                elif str(cell.value).strip():
                    values.append(str(cell.value).strip())
            if values:
                rows.append("\t".join(values))
        sheet_text = f"=== Sheet: {sheet.name} ===\n" + "\n".join(rows)
        parts.append(sheet_text)
        sub_artifacts.append(SubArtifact(index=sheet_index, type="sheet", title=sheet.name, content=sheet_text))
    return ExtractedContent(text="\n\n".join(parts), sub_artifacts=sub_artifacts)


def _convert_office_to_pdf(path: Path, *, document_converter_image: str | None = None) -> Path | None:
    """Convert Office input to a temporary PDF using an isolated LibreOffice profile."""
    command = _resolve_document_command("libreoffice", document_converter_image) or _resolve_document_command(
        "soffice", document_converter_image
    )
    if command is None:
        return None
    output_dir = Path(tempfile.mkdtemp(prefix="apex-office-pdf-"))
    profile_dir = Path(tempfile.mkdtemp(prefix="apex-lo-profile-"))
    output_pdf: Path | None = None
    try:
        process = subprocess.run(
            [
                *command,
                "--headless",
                f"-env:UserInstallation=file://{profile_dir.as_posix()}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(path),
            ],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,
        )
        candidate = output_dir / f"{path.stem}.pdf"
        if process.returncode == 0 and candidate.exists():
            output_pdf = candidate
    except subprocess.TimeoutExpired:
        pass
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)
        if output_pdf is None:
            shutil.rmtree(output_dir, ignore_errors=True)
    return output_pdf


def _pdf_visual_blocks(
    display_name: str,
    path: Path,
    *,
    converted: bool = False,
    document_converter_image: str | None = None,
    page_indices: set[int] | None = None,
) -> list[dict[str, Any]]:
    pages = _pdf_page_data_urls(
        path,
        document_converter_image=document_converter_image,
        page_indices=page_indices,
    )
    label = f"\n{display_name} (rendered from Office):" if converted else f"\n{display_name}:"
    blocks: list[dict[str, Any]] = [{"type": "text", "text": label}]
    for page_number, url in pages:
        blocks.append({"type": "text", "text": f"Page {page_number}"})
        blocks.append({"type": "image_url", "image_url": {"url": url}})
    return blocks


def visual_content_blocks(
    output_dir: Path,
    *,
    document_converter_image: str | None = None,
    pdf_page_indices: dict[str, set[int]] | None = None,
) -> list[dict[str, Any]]:
    """Render artifacts as text and real image blocks; PDFs are rasterized page by page."""
    if not output_dir.is_dir():
        return []
    blocks: list[dict[str, Any]] = []
    converted_pdfs: list[Path] = []
    try:
        for path in sorted(output_dir.iterdir()):
            if not path.is_file():
                continue
            extension = path.suffix.lower()
            try:
                if extension in _TEXT_EXTENSIONS or extension in {".csv", ".ipynb"}:
                    text = extract_file_text(path)
                    if text:
                        blocks.append({"type": "text", "text": f"\n{path.name}:\n{text}"})
                elif extension in _OFFICE_EXTENSIONS:
                    pdf_path = _convert_office_to_pdf(
                        path,
                        document_converter_image=document_converter_image,
                    )
                    if pdf_path is None:
                        text = extract_file_text(path)
                        if text:
                            blocks.append({"type": "text", "text": f"\n{path.name} (text fallback):\n{text}"})
                        continue
                    converted_pdfs.append(pdf_path)
                    blocks.extend(
                        _pdf_visual_blocks(
                            path.name,
                            pdf_path,
                            converted=True,
                            document_converter_image=document_converter_image,
                        )
                    )
                elif extension == ".pdf":
                    blocks.extend(
                        _pdf_visual_blocks(
                            path.name,
                            path,
                            document_converter_image=document_converter_image,
                            page_indices=(pdf_page_indices or {}).get(path.name),
                        )
                    )
                elif mime_type := _IMAGE_MIME_TYPES.get(extension):
                    encoded = base64.b64encode(path.read_bytes()).decode("ascii")
                    blocks.extend(
                        [
                            {"type": "text", "text": f"\n{path.name}:"},
                            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{encoded}"}},
                        ]
                    )
            except Exception as exc:
                blocks.append({"type": "text", "text": f"\n{path.name}: [Error: {exc}]"})
    finally:
        for path in converted_pdfs:
            shutil.rmtree(path.parent, ignore_errors=True)
    return blocks
