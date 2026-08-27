# SPDX-FileCopyrightText: Copyright (c) 2026 NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0

import sqlite3
import sys
import zipfile
from pathlib import Path
from types import SimpleNamespace

import pytest
from openpyxl import Workbook
from pptx import Presentation

from resources_servers.apex_agents.artifacts import (
    artifact_changes_text,
    artifact_text,
    safe_extract_snapshot,
    snapshot_changes,
)
from resources_servers.apex_agents.file_extraction import ExtractedContent, SubArtifact


def _archive(path: Path, files: dict[str, str]) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        for name, value in files.items():
            archive.writestr(name, value)


def test_extracts_and_labels_recursive_artifacts(tmp_path: Path) -> None:
    archive = tmp_path / "snapshot.zip"
    output = tmp_path / "out"
    output.mkdir()
    _archive(archive, {"filesystem/reports/final.txt": "answer", ".apps_data/mail/state.json": "{}"})

    files = safe_extract_snapshot(archive, output, max_files=10, max_uncompressed_bytes=1000)
    text, names = artifact_text(output, files, max_total_chars=1000, max_file_chars=1000)

    assert names == [".apps_data/mail/state.json", "filesystem/reports/final.txt"]
    assert "=== filesystem/reports/final.txt ===\nanswer" in text


def test_rejects_zip_slip(tmp_path: Path) -> None:
    archive = tmp_path / "snapshot.zip"
    output = tmp_path / "out"
    output.mkdir()
    _archive(archive, {"filesystem/../../rubric.txt": "leak"})

    with pytest.raises(ValueError, match="unsafe snapshot path"):
        safe_extract_snapshot(archive, output, max_files=10, max_uncompressed_bytes=1000)


def test_rejects_unexpected_root(tmp_path: Path) -> None:
    archive = tmp_path / "snapshot.zip"
    output = tmp_path / "out"
    output.mkdir()
    _archive(archive, {"grader/rubric.txt": "leak"})

    with pytest.raises(ValueError, match="unexpected snapshot root"):
        safe_extract_snapshot(archive, output, max_files=10, max_uncompressed_bytes=1000)


def test_renders_sqlite_application_state(tmp_path: Path) -> None:
    database = tmp_path / "mail.db"
    connection = sqlite3.connect(database)
    connection.execute("CREATE TABLE messages (recipient TEXT, subject TEXT)")
    connection.execute("INSERT INTO messages VALUES ('person@example.com', 'Completed')")
    connection.commit()
    connection.close()
    root = tmp_path / "root"
    target = root / ".apps_data" / "mail" / "mail.db"
    target.parent.mkdir(parents=True)
    target.write_bytes(database.read_bytes())

    rendered, paths = artifact_text(root, [target], max_total_chars=10_000, max_file_chars=10_000)

    assert paths == [".apps_data/mail/mail.db"]
    assert "Table: messages" in rendered
    assert "person@example.com | Completed" in rendered


def test_extracts_docx_without_gdpval_agent_dependency(monkeypatch, tmp_path: Path) -> None:
    root = tmp_path / "root"
    document_path = root / "filesystem" / "report.docx"
    document_path.parent.mkdir(parents=True)
    document_path.write_bytes(b"docx-placeholder")
    document = SimpleNamespace(paragraphs=[SimpleNamespace(text="APEX-local extraction")])
    monkeypatch.setitem(sys.modules, "docx", SimpleNamespace(Document=lambda _path: document))

    rendered, paths = artifact_text(root, [document_path], max_total_chars=10_000, max_file_chars=10_000)

    assert paths == ["filesystem/report.docx"]
    assert "APEX-local extraction" in rendered


def test_snapshot_changes_excludes_unchanged_and_uses_archipelago_content_tags(tmp_path: Path) -> None:
    initial_archive = tmp_path / "initial.zip"
    final_archive = tmp_path / "final.zip"
    _archive(
        initial_archive,
        {
            "filesystem/unchanged.txt": "same",
            "filesystem/modified.txt": "before",
            "filesystem/deleted.txt": "gone",
        },
    )
    _archive(
        final_archive,
        {
            "filesystem/unchanged.txt": "same",
            "filesystem/modified.txt": "after",
            "filesystem/added.txt": "new",
        },
    )
    initial_root = tmp_path / "initial"
    final_root = tmp_path / "final"
    initial_root.mkdir()
    final_root.mkdir()
    initial_files = safe_extract_snapshot(initial_archive, initial_root, max_files=10, max_uncompressed_bytes=1000)
    final_files = safe_extract_snapshot(final_archive, final_root, max_files=10, max_uncompressed_bytes=1000)

    changes = snapshot_changes(initial_root, initial_files, final_root, final_files)
    rendered = artifact_changes_text(changes, max_total_chars=10_000, max_file_chars=2000)

    assert [(change.path, change.change_type) for change in changes] == [
        ("filesystem/added.txt", "added"),
        ("filesystem/deleted.txt", "deleted"),
        ("filesystem/modified.txt", "modified"),
    ]
    assert "filesystem/unchanged.txt" not in rendered
    assert "<created_content>\nnew\n</created_content>" in rendered
    assert "<deleted_content>\ngone\n</deleted_content>" in rendered
    assert "<diff>\n--- before/filesystem/modified.txt" in rendered
    assert "<updated_content>\nafter\n</updated_content>" in rendered
    assert "--- BEFORE ---" not in rendered


def _add_slide(presentation: Presentation, title: str, body: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body


def test_snapshot_changes_matches_shifted_slides_by_content(tmp_path: Path) -> None:
    initial_root = tmp_path / "initial"
    final_root = tmp_path / "final"
    initial_path = initial_root / "filesystem" / "deck.pptx"
    final_path = final_root / "filesystem" / "deck.pptx"
    initial_path.parent.mkdir(parents=True)
    final_path.parent.mkdir(parents=True)

    initial = Presentation()
    _add_slide(initial, "Alpha", "Stable alpha content")
    _add_slide(initial, "Beta", "Beta result is 10 percent")
    initial.save(initial_path)

    final = Presentation()
    _add_slide(final, "New", "A completely new opening slide")
    _add_slide(final, "Alpha", "Stable alpha content")
    _add_slide(final, "Beta", "Beta result is 12 percent")
    final.save(final_path)

    changes = snapshot_changes(initial_root, [initial_path], final_root, [final_path])

    assert [(change.change_type, change.artifact_type, change.index, change.original_index) for change in changes] == [
        ("modified", "slide", 2, 1),
        ("added", "slide", 0, None),
    ]
    assert changes[0].title == "Beta"
    assert "Beta result is 12 percent" in (changes[0].new_content or "")
    assert changes[1].title == "New"


def test_snapshot_changes_matches_sheets_by_title(tmp_path: Path) -> None:
    initial_root = tmp_path / "initial"
    final_root = tmp_path / "final"
    initial_path = initial_root / "filesystem" / "book.xlsx"
    final_path = final_root / "filesystem" / "book.xlsx"
    initial_path.parent.mkdir(parents=True)
    final_path.parent.mkdir(parents=True)

    initial = Workbook()
    initial.active.title = "Summary"
    initial.active["A1"] = "Old content with no lexical overlap"
    initial.save(initial_path)

    final = Workbook()
    final.active.title = "Summary"
    final.active["A1"] = "Entirely different replacement"
    final.create_sheet("Details")["A1"] = "New sheet"
    final.save(final_path)

    changes = snapshot_changes(initial_root, [initial_path], final_root, [final_path])

    assert [(change.change_type, change.artifact_type, change.title) for change in changes] == [
        ("modified", "sheet", "Summary"),
        ("added", "sheet", "Details"),
    ]
    rendered = artifact_changes_text(changes, max_total_chars=10_000, max_file_chars=5_000)
    assert "<updated_content>\n=== Sheet: Summary ===\nEntirely different replacement" in rendered
    assert "<created_content>\n=== Sheet: Details ===\nNew sheet" in rendered


def test_snapshot_changes_flattens_and_matches_pdf_pages(monkeypatch, tmp_path: Path) -> None:
    initial_root = tmp_path / "initial"
    final_root = tmp_path / "final"
    initial_path = initial_root / "filesystem" / "report.pdf"
    final_path = final_root / "filesystem" / "report.pdf"
    initial_path.parent.mkdir(parents=True)
    final_path.parent.mkdir(parents=True)
    initial_path.write_bytes(b"original-pdf")
    final_path.write_bytes(b"final-pdf")

    def extracted(path: Path, *, document_converter_image: str | None = None) -> ExtractedContent:
        del document_converter_image
        if path == initial_path:
            values = ["Stable introduction", "Revenue was 10 million", "Stable appendix"]
        else:
            values = [
                "Stable introduction",
                "Newly inserted methodology",
                "Revenue was 12 million",
                "Stable appendix",
            ]
        return ExtractedContent(
            text="\n\n".join(values),
            sub_artifacts=[
                SubArtifact(index=index, type="page", title=f"Page {index + 1}", content=value)
                for index, value in enumerate(values)
            ],
        )

    monkeypatch.setattr("resources_servers.apex_agents.artifacts.extract_file_content", extracted)

    changes = snapshot_changes(initial_root, [initial_path], final_root, [final_path])

    assert [(change.change_type, change.artifact_type, change.index, change.original_index) for change in changes] == [
        ("modified", "page", 2, 1),
        ("added", "page", 1, 0),
    ]
    assert all(change.path == "filesystem/report.pdf" for change in changes)
