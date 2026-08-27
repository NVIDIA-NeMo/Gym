# SPDX-FileCopyrightText: Copyright (c) 2026 NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0

import base64
import io
import json
import zipfile
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from resources_servers.apex_agents.file_extraction import (
    _docx_review_text,
    extract_file_content,
    extract_file_text,
    visual_content_blocks,
)


def test_pptx_extraction_includes_table_cells(tmp_path: Path) -> None:
    path = tmp_path / "deliverable.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(7), Inches(2)).table
    table.cell(0, 0).text = "Country–Region"
    table.cell(0, 1).text = "Avoided Curtailment"
    table.cell(1, 0).text = "Germany–North"
    table.cell(1, 1).text = "228,025.94"
    presentation.save(path)

    text = extract_file_text(path)

    assert "Country–Region\tAvoided Curtailment" in text
    assert "Germany–North\t228,025.94" in text


def test_pptx_extraction_preserves_archipelago_slide_records(tmp_path: Path) -> None:
    path = tmp_path / "deliverable.pptx"
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Summary"
    first.placeholders[1].text = "First slide body"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "Appendix"
    second.placeholders[1].text = "Second slide body"
    presentation.save(path)

    extracted = extract_file_content(path)

    assert [(slide.index, slide.type, slide.title) for slide in extracted.sub_artifacts] == [
        (0, "slide", "Summary"),
        (1, "slide", "Appendix"),
    ]
    assert "First slide body" in extracted.sub_artifacts[0].content
    assert "Second slide body" in extracted.sub_artifacts[1].content


def test_docx_extraction_includes_tables(tmp_path: Path) -> None:
    path = tmp_path / "deliverable.docx"
    document = Document()
    document.add_paragraph("Executive summary")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Finding"
    table.cell(1, 0).text = "Thermal Overload"
    table.cell(1, 1).text = "42.4%"
    document.save(path)

    text = extract_file_text(path)

    assert "Executive summary" in text
    assert "Metric\tFinding" in text
    assert "Thermal Overload\t42.4%" in text


def test_docx_review_extraction_includes_redlines_and_anchored_comments() -> None:
    document_xml = b"""<?xml version="1.0" encoding="UTF-8"?>
    <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body>
      <w:p><w:r><w:t>7.7 Termination</w:t></w:r></w:p>
      <w:p><w:commentRangeStart w:id="0"/><w:r><w:t>30 days</w:t></w:r><w:commentRangeEnd w:id="0"/>
        <w:del w:author="Reviewer"><w:r><w:delText>30 days</w:delText></w:r></w:del>
        <w:ins w:author="Agent"><w:r><w:t>45 days</w:t></w:r></w:ins></w:p>
    </w:body></w:document>"""
    comments_xml = b"""<?xml version="1.0" encoding="UTF-8"?>
    <w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
      <w:comment w:id="0" w:author="Agent"><w:p><w:r><w:t>Consider a longer cure period.</w:t></w:r></w:p></w:comment>
    </w:comments>"""
    source = io.BytesIO()
    with zipfile.ZipFile(source, "w") as archive:
        archive.writestr("word/document.xml", document_xml)
        archive.writestr("word/comments.xml", comments_xml)

    text = _docx_review_text(source.getvalue())

    assert '[7.7 Termination] DELETED: "30 days"' in text
    assert '[7.7 Termination] INSERTED: "45 days"' in text
    assert '[7.7 Termination] on "30 days": Consider a longer cure period.' in text


def test_xlsx_extraction_includes_formulas_formats_and_skips_hidden_sheets(tmp_path: Path) -> None:
    path = tmp_path / "deliverable.xlsx"
    workbook = Workbook()
    visible = workbook.active
    visible.title = "Summary"
    visible["A1"] = "Revenue"
    visible["B1"] = 125000
    visible["B1"].number_format = "$#,##0"
    visible["B2"] = "=B1*2"
    hidden = workbook.create_sheet("Raw")
    hidden["A1"] = "should not be graded"
    hidden.sheet_state = "hidden"
    workbook.save(path)

    text = extract_file_text(path)

    assert "=== Sheet: Summary ===" in text
    assert "B1: [fmt: $#,##0]" in text
    assert "B2: =B1*2" in text
    assert "formula cells could not be recalculated" in text
    assert "should not be graded" not in text

    extracted = extract_file_content(path)
    assert [(sheet.index, sheet.type, sheet.title) for sheet in extracted.sub_artifacts] == [(0, "sheet", "Summary")]
    assert "=== Sheet: Summary ===" in extracted.sub_artifacts[0].content
    assert "formula cells could not be recalculated" in extracted.sub_artifacts[0].content


def test_xlsx_extraction_uses_libreoffice_recalculated_values(monkeypatch, tmp_path: Path) -> None:
    path = tmp_path / "formula.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = 125000
    workbook.active["A2"] = "=A1*2"
    workbook.save(path)

    recalculated_path = tmp_path / "recalculated.xlsx"
    recalculated = Workbook()
    recalculated.active["A1"] = 125000
    recalculated.active["A2"] = 250000
    recalculated.save(recalculated_path)
    calls: list[tuple[Path, str | None]] = []

    def fake_recalculate(actual_path: Path, *, document_converter_image: str | None) -> bytes:
        calls.append((actual_path, document_converter_image))
        return recalculated_path.read_bytes()

    monkeypatch.setattr(
        "resources_servers.apex_agents.file_extraction._evaluate_excel_formulas_with_libreoffice",
        fake_recalculate,
    )

    extracted = extract_file_content(path, document_converter_image="/images/archipelago.sif")

    assert calls == [(path, "/images/archipelago.sif")]
    assert "250000" in extracted.text
    assert "A2: =A1*2" in extracted.text
    assert "formula cells could not be recalculated" not in extracted.text


def test_xlsx_extraction_attaches_archipelago_chart_metadata(monkeypatch, tmp_path: Path) -> None:
    path = tmp_path / "chart.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["Quarter", "Revenue"])
    sheet.append(["Q1", 10])
    sheet.append(["Q2", 15])
    chart = BarChart()
    chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=3), titles_from_data=True)
    sheet.add_chart(chart, "D2")
    workbook.save(path)

    pdf_dir = tmp_path / "rendered"
    pdf_dir.mkdir()
    pdf_path = pdf_dir / "chart.pdf"
    pdf_path.write_bytes(b"rendered-pdf")
    monkeypatch.setattr(
        "resources_servers.apex_agents.file_extraction._convert_office_to_pdf",
        lambda *_args, **_kwargs: pdf_path,
    )
    monkeypatch.setattr(
        "resources_servers.apex_agents.file_extraction._pdf_to_chart_images",
        lambda *_args, **_kwargs: [
            {
                "url": "data:image/png;base64,Y2hhcnQ=",
                "placeholder": "[CHART_1]",
                "type": "Chart",
                "caption": "Chart from Excel (Page 1)",
            }
        ],
    )

    extracted = extract_file_content(path, document_converter_image="/images/archipelago.sif")

    assert extracted.images[0]["placeholder"] == "[CHART_1]"
    assert "=== Charts ===\n[CHART_1] - Chart from Excel (Page 1)" in extracted.text


def test_pdf_extraction_marks_image_only_documents(tmp_path: Path) -> None:
    path = tmp_path / "scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with path.open("wb") as stream:
        writer.write(stream)

    text = extract_file_text(path)

    assert "likely scanned/image-only" in text
    assert "=== Page 1 ===" in text
    extracted = extract_file_content(path)
    assert [(page.index, page.type, page.title) for page in extracted.sub_artifacts] == [(0, "page", "Page 1")]
    assert "likely scanned/image-only" in extracted.sub_artifacts[0].content


def test_notebook_extraction_drops_base64_images(tmp_path: Path) -> None:
    path = tmp_path / "analysis.ipynb"
    path.write_text(
        json.dumps(
            {
                "cells": [
                    {"cell_type": "markdown", "source": ["# Findings"]},
                    {
                        "cell_type": "code",
                        "source": ["print('ok')"],
                        "outputs": [
                            {
                                "output_type": "display_data",
                                "data": {"text/plain": ["<Figure>"], "image/png": "very-large-base64"},
                            }
                        ],
                    },
                ]
            }
        )
    )

    text = extract_file_text(path)

    assert "# Findings" in text
    assert "print('ok')" in text
    assert "<Figure>" in text
    assert "very-large-base64" not in text


def test_pdf_visuals_are_png_pages_not_pdf_image_urls(tmp_path: Path, monkeypatch) -> None:
    path = tmp_path / "deliverable.pdf"
    path.write_bytes(b"not parsed because rendering is mocked")
    monkeypatch.setattr("pdf2image.convert_from_path", lambda *args, **kwargs: [Image.new("RGB", (2, 2))])

    blocks = visual_content_blocks(tmp_path)

    image_blocks = [block for block in blocks if block["type"] == "image_url"]
    assert len(image_blocks) == 1
    url = image_blocks[0]["image_url"]["url"]
    assert url.startswith("data:image/png;base64,")
    assert base64.b64decode(url.split(",", 1)[1]).startswith(b"\x89PNG")


def test_pdf_visuals_render_only_requested_page(tmp_path: Path, monkeypatch) -> None:
    path = tmp_path / "deliverable.pdf"
    path.write_bytes(b"not parsed because rendering is mocked")
    calls: list[tuple[int | None, int | None]] = []

    def fake_convert(*_args, **kwargs):
        calls.append((kwargs.get("first_page"), kwargs.get("last_page")))
        return [Image.new("RGB", (2, 2))]

    monkeypatch.setattr("pdf2image.convert_from_path", fake_convert)

    blocks = visual_content_blocks(tmp_path, pdf_page_indices={path.name: {2}})

    assert calls == [(3, 3)]
    assert any(block == {"type": "text", "text": "Page 3"} for block in blocks)
    assert sum(block["type"] == "image_url" for block in blocks) == 1
