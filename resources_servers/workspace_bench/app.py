# SPDX-FileCopyrightText: Copyright (c) 2026 NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0

import asyncio
import csv
import json
import tarfile
import tempfile
from pathlib import Path
from typing import Any

from fastapi import Request
from openai import OpenAI
from pydantic import ConfigDict

from nemo_gym.base_resources_server import (
    BaseResourcesServerConfig,
    BaseSeedSessionRequest,
    BaseSeedSessionResponse,
    BaseVerifyRequest,
    BaseVerifyResponse,
    SimpleResourcesServer,
)
from nemo_gym.global_config import get_global_config_dict
from nemo_gym.sandbox import AsyncSandbox, SandboxResources, SandboxSpec, create_provider
from nemo_gym.sandbox.config import resolve_provider_config, resolve_provider_metadata
from nemo_gym.server_utils import SESSION_ID_KEY, is_nemo_gym_fastapi_entrypoint


class WorkspaceBenchConfig(BaseResourcesServerConfig):
    judge_base_url: str
    judge_api_key: str
    judge_model: str
    sandbox_provider: str
    sandbox_config: dict[str, Any]


class WorkspaceBenchRequest(BaseSeedSessionRequest):
    model_config = ConfigDict(extra="allow")
    task_id: str
    task_dir: str


class WorkspaceBenchSeedResponse(BaseSeedSessionResponse):
    sandbox_handle: str


class WorkspaceBenchVerifyRequest(BaseVerifyRequest):
    model_config = ConfigDict(extra="allow")
    task_id: str
    task_dir: str


class WorkspaceBenchVerifyResponse(BaseVerifyResponse):
    model_config = ConfigDict(extra="allow")
    task_id: str
    passed_count: int
    total_count: int
    judge_model: str
    rubrics: list[dict[str, Any]]


def _extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            from pypdf import PdfReader

            return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
        if suffix == ".docx":
            from docx import Document

            return "\n".join(paragraph.text for paragraph in Document(path).paragraphs)
        if suffix == ".pptx":
            from pptx import Presentation

            return "\n".join(
                shape.text for slide in Presentation(path).slides for shape in slide.shapes if hasattr(shape, "text")
            )
        if suffix == ".xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(path, read_only=True, data_only=True)
            return "\n".join(
                "\t".join(str(value or "") for value in row) for sheet in workbook for row in sheet.values
            )
        if suffix == ".csv":
            return "\n".join("\t".join(row) for row in csv.reader(path.open(encoding="utf-8", errors="replace")))
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as error:
        return f"[content unavailable: {type(error).__name__}]"


def _directory_snapshot(directory: Path, max_bytes: int = 25_000) -> str:
    files = [path for path in sorted(directory.rglob("*")) if path.is_file()]
    if not files:
        return ""
    section_bytes = max_bytes // len(files)
    sections = []
    for path in files:
        header = f"## {path.relative_to(directory)}\n".encode()
        content = _extract_text(path).encode("utf-8")[: max(0, section_bytes - len(header))]
        sections.append(header + content)
    return b"\n\n".join(sections)[:max_bytes].decode("utf-8", errors="ignore")


class WorkspaceBenchResourcesServer(SimpleResourcesServer):
    config: WorkspaceBenchConfig

    def model_post_init(self, context: Any, /) -> None:
        self._sandboxes: dict[str, AsyncSandbox] = {}

    async def seed_session(self, request: Request, body: WorkspaceBenchRequest) -> WorkspaceBenchSeedResponse:
        task_dir = Path(body.task_dir)
        metadata = json.loads((task_dir / "metadata.json").read_text(encoding="utf-8"))
        sandbox_config = self.config.sandbox_config
        provider = create_provider(resolve_provider_config(self.config.sandbox_provider, get_global_config_dict()))
        sandbox = AsyncSandbox(provider)
        spec = SandboxSpec(
            image=sandbox_config["image"],
            ttl_s=sandbox_config.get("ttl_s"),
            ready_timeout_s=sandbox_config.get("ready_timeout_s"),
            workdir="/workspace",
            env={},
            files={},
            metadata=resolve_provider_metadata(self.config.sandbox_provider, get_global_config_dict())
            | sandbox_config.get("metadata", {})
            | {"task-id": body.task_id[:63]},
            resources=SandboxResources.from_mapping(sandbox_config.get("resources", {})),
            provider_options=sandbox_config.get("provider_options", {}),
        )
        await sandbox.start(spec)
        try:
            with tempfile.TemporaryDirectory() as temporary_dir:
                archive = Path(temporary_dir) / "input.tar.gz"
                manifest = metadata.get("data_manifest") or []
                with tarfile.open(archive, "w:gz", dereference=True) as tar:
                    for item in manifest:
                        source = task_dir / item["stored_relpath"]
                        if source.is_file():
                            tar.add(source, arcname=item["filename"])
                await sandbox.upload(archive, "/tmp/input.tar.gz")
            result = await sandbox.exec(
                "mkdir -p /workspace/input /workspace/output /workspace/.opencode "
                "&& ln -s /opt/workspace-bench/office-skills /workspace/.opencode/skills "
                "&& tar -xzf /tmp/input.tar.gz -C /workspace/input",
                cwd="/",
            )
            if result.return_code != 0:
                raise RuntimeError(f"Failed to seed Workspace-Bench input: {result.stderr}")
        except Exception:
            await sandbox.stop()
            raise
        self._sandboxes[str(request.session[SESSION_ID_KEY])] = sandbox
        descriptor = await sandbox.serialize()
        return WorkspaceBenchSeedResponse(sandbox_handle=descriptor["sandbox_id"])

    def _judge(self, metadata: dict[str, Any], input_dir: Path, output_dir: Path) -> list[dict[str, Any]]:
        rubrics = metadata["rubrics"]
        if not any(path.is_file() for path in output_dir.rglob("*")):
            return [
                {
                    "index": index,
                    "rubric": rubric,
                    "passed": False,
                    "confidence": 1.0,
                    "evidence": "No output files found.",
                }
                for index, rubric in enumerate(rubrics)
            ]
        payload = {
            "task": metadata["task"],
            "rubrics": [{"index": index, "rubric": rubric} for index, rubric in enumerate(rubrics)],
            "input_files": _directory_snapshot(input_dir, max_bytes=5_000),
            "output_files": _directory_snapshot(output_dir),
        }
        response = OpenAI(
            base_url=self.config.judge_base_url, api_key=self.config.judge_api_key
        ).chat.completions.create(
            model=self.config.judge_model,
            temperature=0,
            response_format={"type": "json_object"},
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a strict Workspace-Bench evaluator. For every rubric, return a JSON object with a "
                        "rubrics array. Each item must contain index, passed (boolean), confidence (0 to 1), and concise "
                        "file-grounded evidence. Do not award credit without evidence."
                    ),
                },
                {"role": "user", "content": json.dumps(payload, ensure_ascii=False)},
            ],
        )
        judged = json.loads(response.choices[0].message.content or "{}").get("rubrics", [])
        by_index = {item.get("index"): item for item in judged if isinstance(item, dict)}
        return [
            {
                "index": index,
                "rubric": rubric,
                "passed": bool(by_index.get(index, {}).get("passed", False)),
                "confidence": float(by_index.get(index, {}).get("confidence", 0.0)),
                "evidence": str(by_index.get(index, {}).get("evidence", "Judge returned no result.")),
            }
            for index, rubric in enumerate(rubrics)
        ]

    async def verify(self, request: Request, body: WorkspaceBenchVerifyRequest) -> WorkspaceBenchVerifyResponse:
        sandbox = self._sandboxes.pop(str(request.session[SESSION_ID_KEY]))
        try:
            with tempfile.TemporaryDirectory() as temporary_dir:
                local_dir = Path(temporary_dir)
                archive = local_dir / "workspace.tar.gz"
                result = await sandbox.exec("tar -czf /tmp/workspace.tar.gz -C /workspace input output")
                if result.return_code != 0:
                    raise RuntimeError(f"Failed to collect Workspace-Bench files: {result.stderr}")
                await sandbox.download("/tmp/workspace.tar.gz", archive)
                with tarfile.open(archive, "r:gz") as tar:
                    tar.extractall(local_dir, filter="data")
                metadata = json.loads((Path(body.task_dir) / "metadata.json").read_text(encoding="utf-8"))
                rubrics = await asyncio.to_thread(self._judge, metadata, local_dir / "input", local_dir / "output")
        finally:
            await sandbox.stop()
        passed = sum(item["passed"] for item in rubrics)
        total = len(rubrics)
        return WorkspaceBenchVerifyResponse(
            **body.model_dump(),
            reward=passed / total if total else 0.0,
            passed_count=passed,
            total_count=total,
            judge_model=self.config.judge_model,
            rubrics=rubrics,
        )


if __name__ == "__main__":
    WorkspaceBenchResourcesServer.run_webserver()
elif is_nemo_gym_fastapi_entrypoint(__file__):
    app = WorkspaceBenchResourcesServer.run_webserver()  # noqa: F401
